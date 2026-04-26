import collections
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_jobs_style_ppt():
    prs = Presentation()

    # 16:9 Slide size
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Custom color palette
    bg_color = RGBColor(0, 0, 0) # Black background
    text_color = RGBColor(255, 255, 255) # White text
    accent_color = RGBColor(180, 180, 180) # Gray accent

    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
        tf = txBox.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = 'Helvetica Neue'
        run.font.size = Pt(80)
        run.font.bold = True
        run.font.color.rgb = text_color

        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = subtitle
            run2.font.name = 'Helvetica Neue'
            run2.font.size = Pt(40)
            run2.font.color.rgb = accent_color

        return slide

    def add_statement_slide(text):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        txBox = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(3))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = 'Helvetica Neue'
        run.font.size = Pt(65)
        run.font.bold = True
        run.font.color.rgb = text_color

        return slide

    # Slide 1: Title
    add_title_slide("온디바이스 소형 언어모델의 이해", "클라우드 AI에서 개인 기기 AI로")

    # Slide 2: Problem
    add_statement_slide("클라우드는 훌륭합니다.\n하지만 완벽하지 않습니다.")

    # Slide 3: Pain points
    add_statement_slide("비용. 지연 시간. 개인정보.")

    # Slide 4: Solution
    add_statement_slide("해답은 우리의 손 안에 있습니다.\n온디바이스 AI.")

    # Slide 5: Transition to optimization
    add_statement_slide("어떻게 가능한가?")

    # Slide 6: Quantization
    add_statement_slide("1. 모델을 극단적으로 가볍게.\n(양자화 & 1.58-bit)")

    # Slide 7: KV Cache
    add_statement_slide("2. 기억(메모리)의 압축.\n(KV-cache 최소화)")

    # Slide 8: Hardware
    add_statement_slide("3. 맞춤형 심장.\n(모바일 GP-NPU 설계)")

    # Slide 9: Privacy
    add_statement_slide("모든 데이터는\n기기 안에 머뭅니다.")

    # Slide 10: Swarm / Collaboration
    add_statement_slide("스마트폰. 태블릿. 노트북.\n하나의 지능으로 연결됩니다.")

    # Slide 11: Conclusion
    add_statement_slide("더 빠르고, 더 안전하며,\n가장 개인적인 AI.")

    # Slide 12: End
    add_title_slide("감사합니다.", "김현우, 김진희, 김민솔")

    prs.save('presentation.pptx')

if __name__ == '__main__':
    create_jobs_style_ppt()
